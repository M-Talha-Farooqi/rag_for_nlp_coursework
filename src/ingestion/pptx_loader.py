from typing import List
from langchain_core.documents import Document
from pptx import Presentation

class RobustPPTXLoader:
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        """
        Iterates through slides and extracts text from all shapes.
        """
        docs = []
        try:
            prs = Presentation(self.file_path)
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                # Iterate through all shapes (text boxes, titles, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        clean_text = shape.text.strip()
                        if clean_text:
                            slide_text.append(clean_text)
                
                # If we found text on this slide, create a Document
                if slide_text:
                    full_text = "\n".join(slide_text)
                    metadata = {
                        "source": self.file_path,
                        "page": i + 1,  # Slide number
                        "source_filename": self.file_path.split("/")[-1]
                    }
                    docs.append(Document(page_content=full_text, metadata=metadata))
                    
        except Exception as e:
            print(f"⚠️ Error reading PPTX {self.file_path}: {e}")
            
        return docs